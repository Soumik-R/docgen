"""Document file generation service for creating actual .docx and .pptx files."""

from __future__ import annotations

import os
from pathlib import Path
from typing import Optional
from datetime import datetime

from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt


class DocumentGeneratorService:
    """Service to generate actual Word and PowerPoint files."""
    
    def __init__(self, output_dir: str = "generated_documents"):
        """Initialize the document generator.
        
        Args:
            output_dir: Directory to store generated documents
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
    
    def generate_word_document(
        self,
        *,
        title: str,
        content: str,
        document_id: str,
    ) -> str:
        """Generate a Word document (.docx) file.
        
        Args:
            title: Document title
            content: Generated content from AI
            document_id: Unique document identifier
            
        Returns:
            Path to the generated file
        """
        # Create a new Document
        doc = Document()
        
        # Add title
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add metadata
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph(f"Document ID: {document_id}")
        doc.add_paragraph("")  # Empty line
        
        # Add content
        # Split content by paragraphs and add them
        paragraphs = content.split('\n\n')
        for para_text in paragraphs:
            if para_text.strip():
                # Check if it's a heading (starts with #)
                if para_text.strip().startswith('#'):
                    # Remove # symbols and add as heading
                    heading_text = para_text.strip().lstrip('#').strip()
                    level = min(para_text.count('#'), 3)  # Max heading level 3
                    doc.add_heading(heading_text, level=level)
                else:
                    doc.add_paragraph(para_text.strip())
        
        # Save the document
        filename = f"{document_id}_{title.replace(' ', '_')[:30]}.docx"
        filepath = self.output_dir / filename
        doc.save(str(filepath))
        
        return str(filepath)
    
    def generate_powerpoint(
        self,
        *,
        title: str,
        content: str,
        document_id: str,
    ) -> str:
        """Generate a PowerPoint presentation (.pptx) file.
        
        Args:
            title: Presentation title
            content: Generated content from AI
            document_id: Unique document identifier
            
        Returns:
            Path to the generated file
        """
        # Create a new Presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title_shape.text = title
        subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\nDocument ID: {document_id}"
        
        # Split content into slides
        # Each section (separated by double newline or heading) becomes a slide
        sections = content.split('\n\n')
        
        for section in sections:
            if not section.strip():
                continue
                
            # Create a new slide with bullet layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Check if section starts with a heading
            lines = section.strip().split('\n')
            if lines[0].startswith('#'):
                # First line is the slide title
                slide_title = lines[0].lstrip('#').strip()
                slide.shapes.title.text = slide_title
                content_lines = lines[1:]
            else:
                # Use first line as title
                slide.shapes.title.text = lines[0][:50]  # Limit title length
                content_lines = lines[1:] if len(lines) > 1 else [lines[0]]
            
            # Add content to text frame
            if slide.placeholders[1].has_text_frame:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                
                for line in content_lines:
                    if line.strip():
                        p = text_frame.add_paragraph()
                        p.text = line.strip()
                        p.level = 0
        
        # Save the presentation
        filename = f"{document_id}_{title.replace(' ', '_')[:30]}.pptx"
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        
        return str(filepath)
    
    def get_file_path(self, filename: str) -> Optional[Path]:
        """Get the full path to a generated file.
        
        Args:
            filename: Name of the file
            
        Returns:
            Path object if file exists, None otherwise
        """
        filepath = self.output_dir / filename
        return filepath if filepath.exists() else None
